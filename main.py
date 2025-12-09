# =================================================================
# main.py — Universal RAG Backend (FastAPI + Free Stack)
# Features:
# - Multi-file upload (PDF, DOCX, TXT, CSV, XLSX, PPTX, HTML, IMAGES w/ OCR)
# - RAG Q&A with optional filename filter
# - Chat history via conversation_id
# - Streaming responses
# - Image-to-image ASSIST (prompt generator using Ollama)
# - Rate limiting (per IP, in-memory)
# - Caching (answers + summaries, in-memory)
# =================================================================

import os
import uuid
import time
import json
import logging
from typing import List, Optional, Dict, Any, Tuple

import requests
import fitz  # PyMuPDF
import docx
import pytesseract
from PIL import Image, UnidentifiedImageError
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup

from fastapi import FastAPI, UploadFile, File, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse

from qdrant_client import QdrantClient, models
from pydantic import BaseModel, Field
from dotenv import load_dotenv

import nltk
from nltk.tokenize import sent_tokenize

# =================================================================
# Logging
# =================================================================

logger = logging.getLogger("rag_backend")
handler = logging.StreamHandler()
formatter = logging.Formatter(
    "[%(asctime)s] [%(levelname)s] %(name)s - %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
handler.setFormatter(formatter)
if not logger.handlers:
    logger.addHandler(handler)
logger.setLevel(logging.INFO)

# =================================================================
# Environment
# =================================================================

load_dotenv()

# Ollama
OLLAMA_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
LLM_MODEL = os.getenv("LLM_MODEL", "llama3.1")
EMBED_MODEL = os.getenv("EMBEDDING_MODEL", "nomic-embed-text")

# Qdrant
QDRANT_HOST = os.getenv("QDRANT_HOST", "localhost")
QDRANT_PORT = int(os.getenv("QDRANT_PORT", "6333"))
COLLECTION = os.getenv("VECTOR_COLLECTION_NAME", "documents")

# Upload dir
UPLOAD_DIR = os.getenv("UPLOAD_DIR", "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# NLTK
try:
    nltk.data.find("tokenizers/punkt")
except LookupError:
    logger.info("Downloading NLTK punkt...")
    nltk.download("punkt")

# =================================================================
# Qdrant setup
# =================================================================

qdrant = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT)
EMBED_DIM = 768  # for nomic-embed-text

def ensure_collection():
    try:
        qdrant.get_collection(COLLECTION)
        logger.info("Qdrant collection '%s' exists", COLLECTION)
    except Exception:
        logger.info("Creating Qdrant collection '%s' with dim %d", COLLECTION, EMBED_DIM)
        qdrant.create_collection(
            collection_name=COLLECTION,
            vectors_config=models.VectorParams(
                size=EMBED_DIM,
                distance=models.Distance.COSINE
            )
        )

ensure_collection()

# =================================================================
# Rate limiting + caching + chat history
# =================================================================

RATE_LIMIT_REQUESTS = 60      # per IP
RATE_LIMIT_WINDOW = 60.0      # seconds

rate_limit_store: Dict[str, List[float]] = {}  # ip -> [timestamps]

# Simple in-memory caches
CACHE_MAX_ITEMS = 200
answer_cache: Dict[Tuple[str, str], str] = {}   # (filename or ALL, query) -> answer
summary_cache: Dict[str, str] = {}             # filename -> summary

# Chat history: conversation_id -> list[{"role": "user"/"assistant", "content": str}]
chat_histories: Dict[str, List[Dict[str, str]]] = {}

def check_rate_limit(ip: str):
    now = time.time()
    window_start = now - RATE_LIMIT_WINDOW
    timestamps = rate_limit_store.get(ip, [])
    # keep only recent ones
    timestamps = [t for t in timestamps if t >= window_start]
    if len(timestamps) >= RATE_LIMIT_REQUESTS:
        logger.warning("Rate limit exceeded for IP %s", ip)
        raise HTTPException(status_code=429, detail="Too many requests, please slow down.")
    timestamps.append(now)
    rate_limit_store[ip] = timestamps

def cache_put_answer(key: Tuple[str, str], answer: str):
    if len(answer_cache) >= CACHE_MAX_ITEMS:
        # drop one arbitrary item (simple eviction)
        answer_cache.pop(next(iter(answer_cache)))
    answer_cache[key] = answer

def cache_put_summary(filename: str, summary: str):
    if len(summary_cache) >= CACHE_MAX_ITEMS:
        summary_cache.pop(next(iter(summary_cache)))
    summary_cache[filename] = summary

# =================================================================
# Ollama helpers
# =================================================================

def embed_text(text: str) -> List[float]:
    try:
        r = requests.post(
            f"{OLLAMA_URL}/api/embeddings",
            json={"model": EMBED_MODEL, "prompt": text},
            timeout=60,
        )
    except Exception as e:
        logger.exception("Error calling Ollama embeddings")
        raise HTTPException(status_code=502, detail="Embedding service not reachable") from e

    if r.status_code != 200:
        logger.error("Ollama embeddings error: %s", r.text)
        raise HTTPException(status_code=502, detail="Embedding service error")

    data = r.json()
    emb = data.get("embedding")
    if not emb:
        raise HTTPException(status_code=502, detail="Invalid embedding response")
    return emb


def generate_llm(prompt: str) -> str:
    """
    Non-streaming LLM call.
    """
    try:
        r = requests.post(
            f"{OLLAMA_URL}/api/generate",
            json={"model": LLM_MODEL, "prompt": prompt, "stream": False},
            timeout=120,
        )
    except Exception as e:
        logger.exception("Error calling Ollama LLM")
        raise HTTPException(status_code=502, detail="LLM service not reachable") from e

    if r.status_code != 200:
        logger.error("Ollama LLM error: %s", r.text)
        raise HTTPException(status_code=502, detail="LLM service error")

    data = r.json()
    return (data.get("response") or "").strip()


def generate_llm_stream(prompt: str):
    """
    Streaming generator for LLM responses.
    Yields text chunks.
    """
    try:
        r = requests.post(
            f"{OLLAMA_URL}/api/generate",
            json={"model": LLM_MODEL, "prompt": prompt, "stream": True},
            timeout=120,
            stream=True,
        )
    except Exception as e:
        logger.exception("Error calling Ollama LLM (stream)")
        raise HTTPException(status_code=502, detail="LLM service not reachable") from e

    if r.status_code != 200:
        logger.error("Ollama LLM error (stream): %s", r.text)
        raise HTTPException(status_code=502, detail="LLM service error")

    for line in r.iter_lines(decode_unicode=True):
        if not line:
            continue
        try:
            data = json.loads(line)
        except Exception:
            continue
        chunk = data.get("response", "")
        if chunk:
            yield chunk
        if data.get("done"):
            break

# =================================================================
# File extractors
# =================================================================

def extract_pdf(path: str) -> str:
    try:
        text = ""
        pdf = fitz.open(path)
        for p in pdf:
            text += p.get_text()
        return text
    except Exception as e:
        logger.exception("Failed to read PDF: %s", path)
        raise HTTPException(status_code=400, detail="Invalid or corrupted PDF")

def extract_docx(path: str) -> str:
    try:
        d = docx.Document(path)
        return "\n".join(p.text for p in d.paragraphs if p.text.strip())
    except Exception as e:
        logger.exception("Failed to read DOCX: %s", path)
        raise HTTPException(status_code=400, detail="Invalid or corrupted DOCX")

def extract_txt(path: str) -> str:
    for enc in ("utf-8", "latin-1"):
        try:
            with open(path, "r", encoding=enc, errors="ignore") as f:
                return f.read()
        except Exception:
            continue
    return ""

def extract_image(path: str) -> str:
    try:
        img = Image.open(path)
    except UnidentifiedImageError:
        raise HTTPException(status_code=400, detail="Unsupported or corrupted image")
    except Exception:
        raise HTTPException(status_code=400, detail="Failed to open image")

    try:
        text = pytesseract.image_to_string(img)
    except Exception:
        raise HTTPException(
            status_code=500,
            detail="OCR failed. Ensure Tesseract is installed."
        )
    text = (text or "").strip()
    if not text:
        return "No readable text found in the image."
    return text

def extract_csv(path: str) -> str:
    try:
        return pd.read_csv(path).to_string()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid or corrupted CSV")

def extract_excel(path: str) -> str:
    try:
        return pd.read_excel(path).to_string()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid or corrupted Excel")

def extract_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        text = ""
        for s in prs.slides:
            for shape in s.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid or corrupted PPTX")

def extract_html(path: str) -> str:
    try:
        with open(path, "r", errors="ignore") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            return soup.get_text()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid or corrupted HTML")

def extract_any(path: str) -> str:
    try:
        with open(path, "rb") as f:
            return f.read().decode("latin1", errors="ignore")
    except Exception:
        raise HTTPException(status_code=400, detail="Unable to read file")

def extract_file(path: str, filename: str) -> str:
    ext = filename.split(".")[-1].lower()
    if ext == "pdf": return extract_pdf(path)
    if ext == "docx": return extract_docx(path)
    if ext == "txt": return extract_txt(path)
    if ext in ["png", "jpg", "jpeg"]: return extract_image(path)
    if ext == "csv": return extract_csv(path)
    if ext == "xlsx": return extract_excel(path)
    if ext == "pptx": return extract_pptx(path)
    if ext in ["html", "htm"]: return extract_html(path)
    return extract_any(path)

# =================================================================
# Chunking
# =================================================================

def smart_chunk_text(text: str, max_tokens: int = 500) -> List[str]:
    text = (text or "").strip()
    if not text:
        return []

    sentences = sent_tokenize(text)
    chunks: List[str] = []
    current = ""

    for s in sentences:
        if len(current) + len(s) < max_tokens:
            current += " " + s
        else:
            if current.strip():
                chunks.append(current.strip())
            current = s

    if current.strip():
        chunks.append(current.strip())
    return chunks

# =================================================================
# Store + search
# =================================================================

def store_chunk(text: str, filename: str, file_id: str, chunk_id: int) -> None:
    emb = embed_text(text)
    qdrant.upsert(
        collection_name=COLLECTION,
        points=[
            models.PointStruct(
                id=str(uuid.uuid4()),
                vector=emb,
                payload={
                    "text": text,
                    "filename": filename,
                    "file_id": file_id,
                    "chunk_id": chunk_id,
                },
            )
        ],
    )

def search_docs(q_emb: List[float], filename: Optional[str] = None):
    filt = None
    if filename:
        filt = models.Filter(
            must=[models.FieldCondition(key="filename", match=models.MatchValue(value=filename))]
        )

    return qdrant.search(
        collection_name=COLLECTION,
        query_vector=q_emb,
        query_filter=filt,
        limit=5,
    )

# =================================================================
# Pydantic models
# =================================================================

class AskReq(BaseModel):
    query: str = Field(..., min_length=3)
    filename: Optional[str] = None
    conversation_id: Optional[str] = None

class AskResp(BaseModel):
    answer: str
    conversation_id: str
    from_cache: bool = False

class DeleteReq(BaseModel):
    filename: str

class SummarizeReq(BaseModel):
    filename: str

class ImageEditReq(BaseModel):
    instructions: str = Field(..., min_length=3)
    # In future you can include conversation_id, style presets, etc.

# =================================================================
# FastAPI app
# =================================================================

app = FastAPI(title="Universal RAG Backend (Free Stack)")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],   # Restrict in production
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/health")
def health(request: Request):
    check_rate_limit(request.client.host)
    return {"status": "ok"}

# =================================================================
# 1. Upload multiple files
# =================================================================

@app.post("/upload-files")
async def upload_multiple_files(request: Request, files: List[UploadFile] = File(...)):
    check_rate_limit(request.client.host)

    if not files:
        raise HTTPException(status_code=400, detail="No files provided")

    uploaded: List[Dict[str, Any]] = []

    for file in files:
        filename = file.filename or "unnamed"
        logger.info("Processing upload: %s", filename)

        temp_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}_{filename}")
        try:
            with open(temp_path, "wb") as f:
                f.write(await file.read())
        except Exception:
            raise HTTPException(status_code=500, detail="Failed to save uploaded file")

        try:
            extracted = extract_file(temp_path, filename)
        finally:
            try:
                os.remove(temp_path)
            except Exception:
                pass

        chunks = smart_chunk_text(extracted)
        file_id = str(uuid.uuid4())

        for i, c in enumerate(chunks):
            store_chunk(c, filename, file_id, i)

        uploaded.append({
            "filename": filename,
            "file_id": file_id,
            "chunks_stored": len(chunks)
        })

    return {"uploaded": uploaded}

# =================================================================
# 2. List files
# =================================================================

@app.get("/list-files")
def list_files(request: Request):
    check_rate_limit(request.client.host)

    points, _ = qdrant.scroll(collection_name=COLLECTION, limit=100000, with_payload=True)
    counts: Dict[str, int] = {}
    for p in points:
        fn = p.payload.get("filename")
        if not fn:
            continue
        counts[fn] = counts.get(fn, 0) + 1

    return {"files": [{"filename": k, "chunks": v} for k, v in counts.items()]}

# =================================================================
# 3. Delete file
# =================================================================

@app.post("/delete-file")
def delete_file(request: Request, req: DeleteReq):
    check_rate_limit(request.client.host)

    qdrant.delete(
        collection_name=COLLECTION,
        points_selector=models.FilterSelector(
            filter=models.Filter(
                must=[models.FieldCondition(key="filename", match=models.MatchValue(value=req.filename))]
            )
        ),
    )
    return {"status": "deleted", "filename": req.filename}

# =================================================================
# 4. Ask (non-streaming, with history, caching)
# =================================================================

@app.post("/ask", response_model=AskResp)
def ask_question(request: Request, req: AskReq):
    check_rate_limit(request.client.host)

    query = req.query.strip()
    if not query:
        raise HTTPException(status_code=400, detail="Query cannot be empty")

    key = (req.filename or "ALL", query)
    if key in answer_cache:
        return AskResp(
            answer=answer_cache[key],
            conversation_id=req.conversation_id or "cache-only",
            from_cache=True,
        )

    q_emb = embed_text(query)
    hits = search_docs(q_emb, req.filename)

    if not hits:
        answer = "I could not find relevant information in the stored documents."
        cache_put_answer(key, answer)
        return AskResp(answer=answer, conversation_id=req.conversation_id or str(uuid.uuid4()), from_cache=False)

    context_parts = [h.payload.get("text", "") for h in hits]
    context = "\n\n---\n\n".join(context_parts)

    # Manage conversation_id and history
    conv_id = req.conversation_id or str(uuid.uuid4())
    history = chat_histories.get(conv_id, [])

    history.append({"role": "user", "content": query})

    # Build a simple conversational RAG prompt
    history_text = ""
    for msg in history[-6:]:  # last 6 messages
        history_text += f"{msg['role'].upper()}: {msg['content']}\n"

    prompt = f"""
You are a helpful assistant with access to a knowledge base.

CONTEXT (from documents):
{context}

RECENT CHAT HISTORY:
{history_text}

Use ONLY the context above to answer the latest user question.
If the answer is not clearly in the context, say you don't know.

USER QUESTION:
{query}

ASSISTANT ANSWER:
""".strip()

    answer = generate_llm(prompt)
    history.append({"role": "assistant", "content": answer})
    chat_histories[conv_id] = history

    cache_put_answer(key, answer)

    return AskResp(answer=answer, conversation_id=conv_id, from_cache=False)

# =================================================================
# 5. Ask (streaming)
# =================================================================

@app.post("/ask-stream")
def ask_question_stream(request: Request, req: AskReq):
    """
    Streaming version of /ask.
    Returns plain text stream of the answer.
    conversation_id is returned in header X-Conversation-Id.
    """
    check_rate_limit(request.client.host)

    query = req.query.strip()
    if not query:
        raise HTTPException(status_code=400, detail="Query cannot be empty")

    q_emb = embed_text(query)
    hits = search_docs(q_emb, req.filename)

    if not hits:
        def gen_empty():
            yield "I could not find relevant information in the stored documents."
        return StreamingResponse(gen_empty(), media_type="text/plain")

    context_parts = [h.payload.get("text", "") for h in hits]
    context = "\n\n---\n\n".join(context_parts)

    conv_id = req.conversation_id or str(uuid.uuid4())
    history = chat_histories.get(conv_id, [])
    history.append({"role": "user", "content": query})

    history_text = ""
    for msg in history[-6:]:
        history_text += f"{msg['role'].upper()}: {msg['content']}\n"

    prompt = f"""
You are a helpful assistant with access to a knowledge base.

CONTEXT (from documents):
{context}

RECENT CHAT HISTORY:
{history_text}

Use ONLY the context above to answer the latest user question.
If the answer is not clearly in the context, say you don't know.

USER QUESTION:
{query}

ASSISTANT ANSWER:
""".strip()

    # We'll accumulate the full answer to store in history
    def streamer():
        full_answer = ""
        for chunk in generate_llm_stream(prompt):
            full_answer += chunk
            yield chunk
        # after stream finishes, store in history
        history.append({"role": "assistant", "content": full_answer})
        chat_histories[conv_id] = history

    return StreamingResponse(
        streamer(),
        media_type="text/plain",
        headers={"X-Conversation-Id": conv_id},
    )

# =================================================================
# 6. Summarize file (with caching)
# =================================================================

@app.post("/summarize-file")
def summarize_file(request: Request, req: SummarizeReq):
    check_rate_limit(request.client.host)

    filename = req.filename.strip()
    if not filename:
        raise HTTPException(status_code=400, detail="Filename cannot be empty")

    if filename in summary_cache:
        return {"filename": filename, "summary": summary_cache[filename], "from_cache": True}

    filt = models.Filter(
        must=[models.FieldCondition(key="filename", match=models.MatchValue(value=filename))]
    )

    all_texts: List[str] = []
    offset = None
    while True:
        points, next_offset = qdrant.scroll(
            collection_name=COLLECTION,
            limit=256,
            with_payload=True,
            offset=offset,
            scroll_filter=filt,
        )
        for p in points:
            t = p.payload.get("text", "")
            if t:
                all_texts.append(t)
        if next_offset is None:
            break
        offset = next_offset

    if not all_texts:
        raise HTTPException(status_code=404, detail="No chunks found for this filename")

    joined = "\n\n".join(all_texts)
    if len(joined) > 8000:
        joined = joined[:8000]

    prompt = f"""
You are an expert summarizer. Summarize the following document in 5-10 bullet points.
Highlight key concepts, entities, and any important numbers. Be concise and clear.

DOCUMENT CONTENT:
{joined}
""".strip()

    summary = generate_llm(prompt)
    cache_put_summary(filename, summary)

    return {"filename": filename, "summary": summary, "from_cache": False}

# =================================================================
# 7. Image-to-image ASSIST (prompt generator)
# =================================================================

@app.post("/image/assist")
async def image_edit_assist(request: Request, instructions: ImageEditReq, file: UploadFile = File(...)):
    """
    POC endpoint:
    - Takes an image + user instructions
    - Uses Ollama to produce a detailed prompt describing how to transform the image
    - You can feed this prompt into Stable Diffusion / ComfyUI / etc.

    This does NOT actually generate a new image (Ollama is text-only).
    """
    check_rate_limit(request.client.host)

    filename = file.filename or "image"
    temp_path = os.path.join(UPLOAD_DIR, f"img_{uuid.uuid4()}_{filename}")
    try:
        with open(temp_path, "wb") as f:
            f.write(await file.read())
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to save uploaded image")

    # Try to get a brief description via OCR (not perfect but helpful)
    try:
        base_desc = extract_image(temp_path)
    except HTTPException:
        base_desc = "Image content could not be read. Treat it as a generic image."

    try:
        os.remove(temp_path)
    except Exception:
        pass

    prompt = f"""
You are an AI assistant helping to prepare prompts for an image generation / editing model.

The user has an image with the following rough OCR-based description:
'{base_desc}'

They want to apply these instructions:
'{instructions.instructions}'

Produce a SINGLE, detailed, well-structured prompt that can be used for an image generation or
image-to-image model (like Stable Diffusion). Include style, lighting, colors, and other details
if relevant. Do not include technical jargon, just a natural-language prompt.
""".strip()

    detailed_prompt = generate_llm(prompt)
    return {"prompt_for_image_model": detailed_prompt}

# =================================================================
# Run with: uvicorn main:app --reload
# =================================================================
#pip install fastapi uvicorn qdrant-client pymupdf python-docx pytesseract pillow pandas python-pptx beautifulsoup4 requests nltk python-dotenv
